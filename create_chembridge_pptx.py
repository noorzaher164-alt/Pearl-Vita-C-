"""
ChemBridge – Nourhan's Chemistry Tutoring
Professional PowerPoint Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import math
import os

# ── Brand Colors ──────────────────────────────────────────────────────────────
DEEP_PURPLE  = RGBColor(0x3E, 0x1D, 0x68)   # #3E1D68
SOFT_PINK    = RGBColor(0xC9, 0x8B, 0xA6)   # #C98BA6
LIGHT_PURPLE = RGBColor(0x6A, 0x3D, 0x9A)   # mid purple
PALE_PINK    = RGBColor(0xF5, 0xEC, 0xF3)   # near-white pink tint
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x44)
ACCENT_GOLD  = RGBColor(0xE8, 0xD5, 0xB7)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=Pt(0), transparency=0):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if transparency:
            shape.fill.fore_color.theme_color  # dummy access
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(14), font_color=WHITE, bold=False,
                 italic=False, align=PP_ALIGN.LEFT, wrap=True,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox


def add_circle(slide, cx, cy, r, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        9,  # OVAL
        cx - r, cy - r, r * 2, r * 2
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_molecule_decoration(slide, x, y, size=0.3, color=SOFT_PINK, alpha_color=None):
    """Draw a simple molecule: 3 circles connected (decorative)."""
    c = alpha_color or color
    r = Inches(size)
    # central atom
    add_circle(slide, x, y, r, c)
    # arms
    for dx, dy in [(Inches(size*2.2), 0), (-Inches(size*2.2), 0),
                   (0, Inches(size*2.2)), (0, -Inches(size*2.2))]:
        add_circle(slide, x+dx, y+dy, Inches(size*0.7), c)


def add_atom_ring(slide, cx, cy, radius=Inches(0.6), color=SOFT_PINK):
    """Draw a benzene-ring-like decorative atom."""
    # outer circle (ring)
    ring = slide.shapes.add_shape(9, cx - radius, cy - radius, radius*2, radius*2)
    ring.fill.background()
    ring.line.color.rgb = color
    ring.line.width = Pt(1.5)
    # nucleus
    add_circle(slide, cx, cy, Inches(0.12), color)


def gradient_header(slide, height=Inches(1.4)):
    """Two-layer gradient header bar."""
    add_rect(slide, 0, 0, W, height, fill_color=DEEP_PURPLE)
    # accent stripe
    add_rect(slide, 0, height - Inches(0.07), W, Inches(0.07), fill_color=SOFT_PINK)


def slide_footer(slide, text="ChemBridge – Nourhan's Chemistry Tutoring  |  chembridge.msnour.com"):
    footer_h = Inches(0.32)
    add_rect(slide, 0, H - footer_h, W, footer_h, fill_color=DEEP_PURPLE)
    add_text_box(slide, text,
                 Inches(0.3), H - footer_h + Inches(0.04),
                 W - Inches(0.6), footer_h,
                 font_size=Pt(8), font_color=SOFT_PINK,
                 align=PP_ALIGN.CENTER)


def section_tag(slide, text, left, top):
    """Small pill-shaped section label."""
    w, h = Inches(2.2), Inches(0.32)
    r = add_rect(slide, left, top, w, h, fill_color=SOFT_PINK)
    add_text_box(slide, text, left, top + Inches(0.04), w, h,
                 font_size=Pt(9), font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)


def step_box(slide, number, title, body, left, top, w=Inches(2.2), h=Inches(1.3)):
    """Numbered step card."""
    add_rect(slide, left, top, w, h, fill_color=DEEP_PURPLE,
             line_color=SOFT_PINK, line_width=Pt(1))
    # number circle
    add_circle(slide, left + Inches(0.3), top + Inches(0.32), Inches(0.22), SOFT_PINK)
    add_text_box(slide, str(number),
                 left + Inches(0.1), top + Inches(0.1), Inches(0.4), Inches(0.45),
                 font_size=Pt(11), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 left + Inches(0.55), top + Inches(0.06), w - Inches(0.6), Inches(0.35),
                 font_size=Pt(10), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, body,
                 left + Inches(0.12), top + Inches(0.48), w - Inches(0.2), h - Inches(0.55),
                 font_size=Pt(8.5), font_color=WHITE)


def feature_card(slide, icon, title, body, left, top, w=Inches(2.8), h=Inches(1.55)):
    add_rect(slide, left, top, w, h, fill_color=PALE_PINK,
             line_color=SOFT_PINK, line_width=Pt(1.2))
    add_text_box(slide, icon,
                 left + Inches(0.12), top + Inches(0.12), Inches(0.6), Inches(0.55),
                 font_size=Pt(22), font_color=DEEP_PURPLE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 left + Inches(0.75), top + Inches(0.14), w - Inches(0.85), Inches(0.38),
                 font_size=Pt(11), font_color=DEEP_PURPLE, bold=True)
    add_text_box(slide, body,
                 left + Inches(0.12), top + Inches(0.62), w - Inches(0.24), h - Inches(0.72),
                 font_size=Pt(8.5), font_color=DARK_GRAY)


def arrow_down(slide, cx, top, color=SOFT_PINK):
    """Draw a downward arrow connector."""
    h_line = Inches(0.25)
    add_rect(slide, cx - Inches(0.03), top, Inches(0.06), h_line, fill_color=color)
    # arrowhead triangle – use a connector shape
    tri = slide.shapes.add_shape(
        5,  # RIGHT_TRIANGLE – we'll rotate
        cx - Inches(0.12), top + h_line, Inches(0.24), Inches(0.14)
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    tri.rotation = 180


def flow_node(slide, text, cx, top, w=Inches(2.6), h=Inches(0.52), bg=DEEP_PURPLE, fg=WHITE):
    left = cx - w / 2
    add_rect(slide, left, top, w, h, fill_color=bg,
             line_color=SOFT_PINK, line_width=Pt(1))
    add_text_box(slide, text, left + Inches(0.1), top + Inches(0.08),
                 w - Inches(0.2), h - Inches(0.1),
                 font_size=Pt(10), font_color=fg, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – COVER PAGE
# ══════════════════════════════════════════════════════════════════════════════
def slide01():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)

    # Full left panel (dark purple)
    add_rect(slide, 0, 0, W * 0.55, H, fill_color=DEEP_PURPLE)

    # Decorative chemistry elements on left panel
    for cx, cy, r in [
        (Inches(1.2), Inches(1.1), Inches(0.35)),
        (Inches(6.5), Inches(6.8), Inches(0.28)),
        (Inches(0.5), Inches(5.5), Inches(0.22)),
    ]:
        add_atom_ring(slide, cx, cy, r, SOFT_PINK)

    add_molecule_decoration(slide, Inches(6.2), Inches(1.5), 0.18, SOFT_PINK)
    add_molecule_decoration(slide, Inches(1.8), Inches(6.2), 0.16, SOFT_PINK)

    # Pink accent bar at bottom of left panel
    add_rect(slide, 0, H - Inches(0.8), W * 0.55, Inches(0.08), fill_color=SOFT_PINK)

    # ── Right side decorative ──
    add_molecule_decoration(slide, Inches(10.5), Inches(1.2), 0.22, SOFT_PINK)
    add_atom_ring(slide, Inches(12.0), Inches(5.5), Inches(0.45), PALE_PINK)
    add_molecule_decoration(slide, Inches(8.8), Inches(6.5), 0.18, PALE_PINK)

    # ── Left panel content ──
    # Logo placeholder circle
    logo_cx, logo_cy, logo_r = Inches(3.65), Inches(1.6), Inches(0.75)
    add_circle(slide, logo_cx, logo_cy, logo_r, SOFT_PINK)
    add_circle(slide, logo_cx, logo_cy, logo_r - Inches(0.06), WHITE)
    add_text_box(slide, "CB",
                 logo_cx - logo_r, logo_cy - logo_r,
                 logo_r * 2, logo_r * 2,
                 font_size=Pt(28), font_color=DEEP_PURPLE, bold=True,
                 align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, "Welcome to",
                 Inches(0.5), Inches(2.65), Inches(6.6), Inches(0.6),
                 font_size=Pt(20), font_color=SOFT_PINK, italic=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "ChemBridge",
                 Inches(0.5), Inches(3.2), Inches(6.6), Inches(1.0),
                 font_size=Pt(44), font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.8), Inches(4.18), Inches(3.6), Inches(0.04),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Your Complete Chemistry Learning Platform",
                 Inches(0.5), Inches(4.3), Inches(6.6), Inches(0.55),
                 font_size=Pt(13), font_color=SOFT_PINK,
                 align=PP_ALIGN.CENTER)

    # Teacher name
    add_text_box(slide, "by Teacher Nourhan Zaher",
                 Inches(0.5), Inches(5.1), Inches(6.6), Inches(0.42),
                 font_size=Pt(11), font_color=ACCENT_GOLD, italic=True,
                 align=PP_ALIGN.CENTER)

    # URL badge
    badge = add_rect(slide, Inches(1.4), Inches(5.68), Inches(4.2), Inches(0.42),
                     fill_color=SOFT_PINK)
    add_text_box(slide, "chembridge.msnour.com",
                 Inches(1.4), Inches(5.72), Inches(4.2), Inches(0.38),
                 font_size=Pt(10), font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)

    # ── Right panel ──
    add_text_box(slide, "Nourhan's Chemistry Tutoring",
                 Inches(7.5), Inches(2.9), Inches(5.0), Inches(0.55),
                 font_size=Pt(15), font_color=DEEP_PURPLE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "Chemistry for Qatar University\nCHEM 101 & CHEM 103",
                 Inches(7.5), Inches(3.52), Inches(5.0), Inches(0.75),
                 font_size=Pt(12), font_color=LIGHT_PURPLE,
                 align=PP_ALIGN.CENTER)

    # Feature pill badges on right
    for i, label in enumerate(["📹  Video Lessons", "📝  Study Notes",
                                "✅  Quizzes", "📊  Progress Tracking"]):
        by = Inches(4.55) + i * Inches(0.55)
        add_rect(slide, Inches(7.9), by, Inches(4.1), Inches(0.4),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.8))
        add_text_box(slide, label, Inches(7.9), by + Inches(0.06),
                     Inches(4.1), Inches(0.35),
                     font_size=Pt(10), font_color=DEEP_PURPLE, align=PP_ALIGN.CENTER)

    slide_footer(slide)

    # Speaker notes
    slide.notes_slide.notes_text_frame.text = (
        "Welcome everyone! This presentation will walk you through everything you "
        "need to know about using ChemBridge – Nourhan's Chemistry Tutoring platform. "
        "ChemBridge is designed specifically for Qatar University students taking "
        "CHEM 101 and CHEM 103. Let's get started!"
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – ABOUT CHEMBRIDGE
# ══════════════════════════════════════════════════════════════════════════════
def slide02():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)
    add_molecule_decoration(slide, Inches(0.6), Inches(6.8), 0.2, PALE_PINK)

    add_text_box(slide, "ABOUT",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "About ChemBridge",
                 Inches(0.5), Inches(0.52), Inches(8.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Mission banner
    add_rect(slide, Inches(0.5), Inches(1.65), W - Inches(1.0), Inches(0.72),
             fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(1))
    add_text_box(slide,
                 "🎯  ChemBridge is an online learning platform built to help Qatar University students "
                 "master CHEM 101 & CHEM 103 through structured video lessons, practice questions, and smart quizzes.",
                 Inches(0.7), Inches(1.72), W - Inches(1.4), Inches(0.62),
                 font_size=Pt(11), font_color=DEEP_PURPLE, align=PP_ALIGN.CENTER)

    # Feature cards row 1
    cards = [
        ("🎓", "Expert Guidance",
         "Taught by Nourhan Zaher, an experienced chemistry instructor specialized in QU curriculum."),
        ("🏗️", "Structured Curriculum",
         "Content is organized by chapter and topic — aligned exactly with your university syllabus."),
        ("⏰", "Study Anytime",
         "Access lessons, notes, and quizzes 24/7 from any device — at your own pace."),
        ("📈", "Track Your Progress",
         "Monitor quiz scores and completed lessons to stay on top of your study goals."),
    ]
    for i, (icon, title, body) in enumerate(cards):
        lx = Inches(0.5) + i * Inches(3.18)
        feature_card(slide, icon, title, body, lx, Inches(2.55), Inches(3.05), Inches(1.65))

    # Bottom row
    cards2 = [
        ("📹", "HD Video Lessons",
         "Clear, high-quality recorded lectures you can re-watch as many times as needed."),
        ("📝", "Study Materials",
         "Downloadable notes, formula sheets, and chapter summaries at your fingertips."),
        ("✅", "Practice & Quizzes",
         "Test your knowledge with targeted exercises and timed quizzes after each chapter."),
        ("💬", "Direct Support",
         "Reach Teacher Nourhan directly via WhatsApp for any questions or help."),
    ]
    for i, (icon, title, body) in enumerate(cards2):
        lx = Inches(0.5) + i * Inches(3.18)
        feature_card(slide, icon, title, body, lx, Inches(4.38), Inches(3.05), Inches(1.65))

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "ChemBridge was created specifically for Qatar University chemistry students. "
        "The platform brings together everything a student needs: structured video lectures, "
        "comprehensive notes, practice problems, and quizzes — all aligned with the QU CHEM 101 "
        "and CHEM 103 curriculum. Students can study at their own pace and always reach out to "
        "Teacher Nourhan for personal support."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – HOMEPAGE OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
def slide03():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "NAVIGATION",
                 Inches(0.5), Inches(0.22), Inches(2.5), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Homepage Overview",
                 Inches(0.5), Inches(0.52), Inches(9.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Mock browser screenshot area
    browser_l, browser_t = Inches(0.5), Inches(1.55)
    browser_w, browser_h = Inches(8.2), Inches(5.35)
    add_rect(slide, browser_l, browser_t, browser_w, browser_h,
             fill_color=DARK_GRAY, line_color=SOFT_PINK, line_width=Pt(1.5))

    # Browser chrome
    add_rect(slide, browser_l, browser_t, browser_w, Inches(0.36),
             fill_color=RGBColor(0x22, 0x22, 0x33))
    for cx in [Inches(0.7), Inches(0.95), Inches(1.2)]:
        add_circle(slide, cx, browser_t + Inches(0.18), Inches(0.07),
                   RGBColor(0x88, 0x88, 0x88))
    # address bar
    add_rect(slide, Inches(1.5), browser_t + Inches(0.07),
             Inches(5.5), Inches(0.22), fill_color=RGBColor(0x44, 0x44, 0x55))
    add_text_box(slide, "chembridge.msnour.com",
                 Inches(1.55), browser_t + Inches(0.09), Inches(5.4), Inches(0.2),
                 font_size=Pt(8), font_color=SOFT_PINK, align=PP_ALIGN.CENTER)

    # Navbar inside browser
    add_rect(slide, browser_l, browser_t + Inches(0.36), browser_w, Inches(0.42),
             fill_color=DEEP_PURPLE)
    nav_items = ["Home", "Courses", "Notes", "Quizzes", "Progress", "Contact"]
    for i, item in enumerate(nav_items):
        add_text_box(slide, item,
                     Inches(1.0) + i * Inches(1.25), browser_t + Inches(0.42),
                     Inches(1.15), Inches(0.38),
                     font_size=Pt(9), font_color=WHITE, bold=(i == 0),
                     align=PP_ALIGN.CENTER)

    # Hero section mock
    add_rect(slide, browser_l, browser_t + Inches(0.78), browser_w, Inches(1.6),
             fill_color=RGBColor(0x55, 0x30, 0x80))
    add_text_box(slide, "Welcome to ChemBridge",
                 Inches(0.9), browser_t + Inches(0.98), Inches(5.0), Inches(0.55),
                 font_size=Pt(17), font_color=WHITE, bold=True)
    add_text_box(slide, "Master CHEM 101 & CHEM 103 with structured lessons, notes & quizzes",
                 Inches(0.9), browser_t + Inches(1.48), Inches(5.0), Inches(0.45),
                 font_size=Pt(9), font_color=SOFT_PINK)
    add_rect(slide, Inches(0.9), browser_t + Inches(1.98), Inches(1.4), Inches(0.28),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Get Started →",
                 Inches(0.9), browser_t + Inches(2.0), Inches(1.4), Inches(0.26),
                 font_size=Pt(8.5), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.45), browser_t + Inches(1.98), Inches(1.4), Inches(0.28),
             fill_color=RGBColor(0x44, 0x22, 0x66))
    add_text_box(slide, "View Courses",
                 Inches(2.45), browser_t + Inches(2.0), Inches(1.4), Inches(0.26),
                 font_size=Pt(8.5), font_color=SOFT_PINK, bold=True, align=PP_ALIGN.CENTER)

    # Course preview cards inside browser
    for i, course in enumerate(["CHEM 101", "CHEM 103"]):
        cx = Inches(0.9) + i * Inches(3.8)
        add_rect(slide, cx, browser_t + Inches(2.5), Inches(3.5), Inches(1.1),
                 fill_color=RGBColor(0x2A, 0x12, 0x44),
                 line_color=SOFT_PINK, line_width=Pt(0.8))
        add_text_box(slide, course,
                     cx + Inches(0.15), browser_t + Inches(2.6), Inches(2.5), Inches(0.38),
                     font_size=Pt(13), font_color=SOFT_PINK, bold=True)
        add_text_box(slide, "General Chemistry I" if "101" in course else "General Chemistry II",
                     cx + Inches(0.15), browser_t + Inches(2.98), Inches(3.0), Inches(0.28),
                     font_size=Pt(8.5), font_color=WHITE)

    # Annotation callouts on right
    annotations = [
        (Inches(9.05), Inches(1.85), "🔷 Top Navigation Bar\nQuick access to all sections"),
        (Inches(9.05), Inches(2.65), "🔷 Hero Section\nWelcome banner & CTA buttons"),
        (Inches(9.05), Inches(3.5),  "🔷 Course Cards\nCHEM 101 & CHEM 103 shortcuts"),
        (Inches(9.05), Inches(4.35), "🔷 Login / Sign Up\nAccess your personal dashboard"),
    ]
    for lx, ty, txt in annotations:
        add_rect(slide, lx, ty, Inches(4.0), Inches(0.65),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.8))
        add_text_box(slide, txt, lx + Inches(0.1), ty + Inches(0.07),
                     Inches(3.8), Inches(0.58),
                     font_size=Pt(9), font_color=DEEP_PURPLE)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "When you first open ChemBridge, you'll see the homepage with a clean navigation bar "
        "at the top. It has links to Courses, Notes, Quizzes, Progress, and Contact. "
        "The hero section shows a welcome message and two main buttons: Get Started and View Courses. "
        "Below that you'll find course cards for CHEM 101 and CHEM 103 which take you directly "
        "into your relevant course content."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – CREATING AN ACCOUNT
# ══════════════════════════════════════════════════════════════════════════════
def slide04():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "GET STARTED",
                 Inches(0.5), Inches(0.22), Inches(2.5), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "How to Create Your Account",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Important note
    add_rect(slide, Inches(0.5), Inches(1.58), W - Inches(1.0), Inches(0.52),
             fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(1))
    add_text_box(slide,
                 "⚠️  Accounts are created by Teacher Nourhan. You will receive a welcome email "
                 "with a link to set your password — no self-registration required.",
                 Inches(0.7), Inches(1.65), W - Inches(1.4), Inches(0.42),
                 font_size=Pt(10), font_color=DEEP_PURPLE, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Receive Invitation",
         "Check your email for a welcome message from ChemBridge with your account link."),
        ("2", "Open the Link",
         "Click the 'Set my password' button in the email or copy the link into your browser."),
        ("3", "Set Your Password",
         "Choose a strong password. Enter it twice to confirm, then click Save."),
        ("4", "Verify Your Account",
         "Your account is now active. You'll be redirected to the ChemBridge login page."),
        ("5", "Log In",
         "Enter your email and the new password you just set to access your dashboard."),
    ]
    step_w = Inches(2.4)
    step_h = Inches(3.5)
    gap = Inches(0.2)
    total_w = len(steps) * step_w + (len(steps) - 1) * gap
    start_x = (W - total_w) / 2

    for i, (num, title, body) in enumerate(steps):
        lx = start_x + i * (step_w + gap)
        ty = Inches(2.3)
        # card
        add_rect(slide, lx, ty, step_w, step_h,
                 fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1.2))
        # number circle
        add_circle(slide, lx + step_w / 2, ty + Inches(0.5), Inches(0.32), SOFT_PINK)
        add_text_box(slide, num,
                     lx + step_w / 2 - Inches(0.32), ty + Inches(0.2),
                     Inches(0.64), Inches(0.58),
                     font_size=Pt(16), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # title
        add_text_box(slide, title,
                     lx + Inches(0.1), ty + Inches(1.05), step_w - Inches(0.2), Inches(0.5),
                     font_size=Pt(10), font_color=SOFT_PINK, bold=True, align=PP_ALIGN.CENTER)
        # body
        add_text_box(slide, body,
                     lx + Inches(0.15), ty + Inches(1.6), step_w - Inches(0.3), Inches(1.65),
                     font_size=Pt(8.5), font_color=WHITE, align=PP_ALIGN.CENTER)
        # connector arrow (not after last)
        if i < len(steps) - 1:
            arr_x = lx + step_w + gap / 2
            arr_y = ty + step_h / 2 - Inches(0.05)
            add_text_box(slide, "→", arr_x - Inches(0.15), arr_y - Inches(0.18),
                         Inches(0.3), Inches(0.36),
                         font_size=Pt(18), font_color=SOFT_PINK, align=PP_ALIGN.CENTER)

    # Email preview box
    add_rect(slide, Inches(0.5), Inches(6.0), W - Inches(1.0), Inches(0.9),
             fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(1))
    add_text_box(slide, "📧  Email Subject: 'Welcome to ChemBridge - set your password'  |  "
                 "From: ChemBridge <chembridge@msnour.com>  |  Button: Set my password",
                 Inches(0.7), Inches(6.1), W - Inches(1.4), Inches(0.72),
                 font_size=Pt(9.5), font_color=DEEP_PURPLE, align=PP_ALIGN.CENTER)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "Unlike most platforms, ChemBridge accounts are set up by Teacher Nourhan — "
        "students do not register themselves. Once enrolled, you receive a welcome email "
        "from chembridge@msnour.com with a 'Set my password' button. Click it, create your "
        "password, and you're ready to log in. This keeps the platform secure and ensures "
        "only enrolled students have access."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – LOGGING IN
# ══════════════════════════════════════════════════════════════════════════════
def slide05():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "ACCESS",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Logging In to ChemBridge",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Left: login form mockup
    form_l, form_t = Inches(0.8), Inches(1.65)
    form_w, form_h = Inches(5.2), Inches(4.8)
    add_rect(slide, form_l, form_t, form_w, form_h,
             fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1.5))

    # Logo area in form
    add_circle(slide, form_l + form_w / 2, form_t + Inches(0.7), Inches(0.42), SOFT_PINK)
    add_text_box(slide, "CB",
                 form_l + form_w / 2 - Inches(0.42), form_t + Inches(0.3),
                 Inches(0.84), Inches(0.8),
                 font_size=Pt(18), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Sign in to ChemBridge",
                 form_l + Inches(0.4), form_t + Inches(1.35),
                 form_w - Inches(0.8), Inches(0.42),
                 font_size=Pt(13), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Email field
    add_text_box(slide, "Email Address",
                 form_l + Inches(0.4), form_t + Inches(1.95),
                 form_w - Inches(0.8), Inches(0.3),
                 font_size=Pt(9), font_color=SOFT_PINK)
    add_rect(slide, form_l + Inches(0.4), form_t + Inches(2.28),
             form_w - Inches(0.8), Inches(0.38),
             fill_color=RGBColor(0x55, 0x30, 0x80),
             line_color=SOFT_PINK, line_width=Pt(0.8))
    add_text_box(slide, "your.email@qu.edu.qa",
                 form_l + Inches(0.5), form_t + Inches(2.32),
                 form_w - Inches(1.0), Inches(0.3),
                 font_size=Pt(9), font_color=RGBColor(0xAA, 0xAA, 0xBB))

    # Password field
    add_text_box(slide, "Password",
                 form_l + Inches(0.4), form_t + Inches(2.82),
                 form_w - Inches(0.8), Inches(0.3),
                 font_size=Pt(9), font_color=SOFT_PINK)
    add_rect(slide, form_l + Inches(0.4), form_t + Inches(3.15),
             form_w - Inches(0.8), Inches(0.38),
             fill_color=RGBColor(0x55, 0x30, 0x80),
             line_color=SOFT_PINK, line_width=Pt(0.8))
    add_text_box(slide, "••••••••••",
                 form_l + Inches(0.5), form_t + Inches(3.19),
                 form_w - Inches(1.0), Inches(0.3),
                 font_size=Pt(9), font_color=WHITE)

    # Login button
    add_rect(slide, form_l + Inches(0.4), form_t + Inches(3.72),
             form_w - Inches(0.8), Inches(0.44),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Sign In  →",
                 form_l + Inches(0.4), form_t + Inches(3.77),
                 form_w - Inches(0.8), Inches(0.38),
                 font_size=Pt(11), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Right: tips
    tips_l = Inches(6.5)
    add_text_box(slide, "Login Tips",
                 tips_l, Inches(1.7), Inches(6.0), Inches(0.5),
                 font_size=Pt(18), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, tips_l, Inches(2.18), Inches(0.06), Inches(2.0),
             fill_color=SOFT_PINK)

    tip_items = [
        ("📧", "Use the email where you received your invitation"),
        ("🔑", "Enter the password you set via the invitation link"),
        ("🔒", "Your session stays active — no need to log in every time"),
        ("❓", "Forgot password? Contact Teacher Nourhan on WhatsApp"),
        ("📱", "Works on mobile, tablet, and desktop browsers"),
        ("🌐", "Always access via: chembridge.msnour.com"),
    ]
    for i, (icon, tip) in enumerate(tip_items):
        ty = Inches(2.4) + i * Inches(0.65)
        add_circle(slide, tips_l + Inches(0.22), ty + Inches(0.18), Inches(0.18), PALE_PINK)
        add_text_box(slide, icon,
                     tips_l + Inches(0.04), ty + Inches(0.02),
                     Inches(0.36), Inches(0.36),
                     font_size=Pt(12), align=PP_ALIGN.CENTER)
        add_text_box(slide, tip,
                     tips_l + Inches(0.52), ty + Inches(0.04),
                     Inches(5.5), Inches(0.38),
                     font_size=Pt(10), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "Logging in is simple. Go to chembridge.msnour.com and click Sign In. "
        "Enter the email address where you received your invitation and the password you created. "
        "The site works on any device — phone, tablet, or computer. If you forget your password, "
        "reach out to Teacher Nourhan on WhatsApp at 33625942."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – EXPLORING COURSES
# ══════════════════════════════════════════════════════════════════════════════
def slide06():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "COURSES",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Exploring Your Courses",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Course catalog label
    add_text_box(slide, "📚  Course Catalog",
                 Inches(0.5), Inches(1.6), Inches(4.0), Inches(0.42),
                 font_size=Pt(14), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, Inches(0.5), Inches(1.98), Inches(4.0), Inches(0.04),
             fill_color=SOFT_PINK)

    # CHEM 101 Card – big featured
    add_rect(slide, Inches(0.5), Inches(2.12), Inches(5.8), Inches(2.3),
             fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(2))
    add_text_box(slide, "FEATURED",
                 Inches(0.68), Inches(2.22), Inches(1.5), Inches(0.28),
                 font_size=Pt(7.5), font_color=SOFT_PINK, bold=True)
    add_rect(slide, Inches(0.68), Inches(2.22), Inches(1.5), Inches(0.28),
             fill_color=SOFT_PINK)
    add_text_box(slide, "FEATURED",
                 Inches(0.68), Inches(2.24), Inches(1.5), Inches(0.26),
                 font_size=Pt(7.5), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "CHEM 101",
                 Inches(0.68), Inches(2.56), Inches(5.3), Inches(0.65),
                 font_size=Pt(26), font_color=WHITE, bold=True)
    add_text_box(slide, "General Chemistry I",
                 Inches(0.68), Inches(3.18), Inches(5.3), Inches(0.38),
                 font_size=Pt(13), font_color=SOFT_PINK)
    add_text_box(slide, "Atomic structure · Bonding · Thermodynamics · Kinetics · Equilibrium",
                 Inches(0.68), Inches(3.55), Inches(5.3), Inches(0.55),
                 font_size=Pt(9), font_color=WHITE)
    add_rect(slide, Inches(0.68), Inches(4.1), Inches(1.5), Inches(0.32),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Open Course →",
                 Inches(0.68), Inches(4.13), Inches(1.5), Inches(0.28),
                 font_size=Pt(8.5), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # CHEM 103 Card
    add_rect(slide, Inches(0.5), Inches(4.58), Inches(5.8), Inches(1.85),
             fill_color=RGBColor(0x55, 0x30, 0x80), line_color=SOFT_PINK, line_width=Pt(1.2))
    add_text_box(slide, "CHEM 103",
                 Inches(0.68), Inches(4.7), Inches(5.3), Inches(0.55),
                 font_size=Pt(22), font_color=WHITE, bold=True)
    add_text_box(slide, "General Chemistry II",
                 Inches(0.68), Inches(5.22), Inches(5.3), Inches(0.35),
                 font_size=Pt(12), font_color=SOFT_PINK)
    add_text_box(slide, "Electrochemistry · Acids & Bases · Organic Introduction · Solutions",
                 Inches(0.68), Inches(5.55), Inches(5.3), Inches(0.45),
                 font_size=Pt(9), font_color=WHITE)

    # Right side info panel
    info_l = Inches(6.8)
    add_text_box(slide, "What's Inside Each Course?",
                 info_l, Inches(1.6), Inches(6.0), Inches(0.48),
                 font_size=Pt(14), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, info_l, Inches(2.05), Inches(3.5), Inches(0.04),
             fill_color=SOFT_PINK)

    course_features = [
        ("🎬", "Video Lectures", "Chapter-by-chapter recorded lessons"),
        ("📄", "Lecture Notes",  "Printable PDFs aligned with each video"),
        ("📋", "Practice Sets",  "Curated exercises per topic"),
        ("🧪", "Lab Concepts",   "Theoretical lab preparation materials"),
        ("🏆", "Chapter Quizzes","Timed self-assessment quizzes"),
        ("📊", "Progress Board", "See your scores and completion rate"),
    ]
    for i, (icon, title, desc) in enumerate(course_features):
        row = i // 2
        col = i % 2
        bx = info_l + col * Inches(3.0)
        by = Inches(2.22) + row * Inches(1.2)
        add_rect(slide, bx, by, Inches(2.85), Inches(1.0),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.8))
        add_text_box(slide, icon,
                     bx + Inches(0.1), by + Inches(0.1), Inches(0.5), Inches(0.5),
                     font_size=Pt(18))
        add_text_box(slide, title,
                     bx + Inches(0.62), by + Inches(0.1), Inches(2.15), Inches(0.35),
                     font_size=Pt(10), font_color=DEEP_PURPLE, bold=True)
        add_text_box(slide, desc,
                     bx + Inches(0.62), by + Inches(0.48), Inches(2.15), Inches(0.42),
                     font_size=Pt(8.5), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "ChemBridge currently offers two courses: CHEM 101 (General Chemistry I) and CHEM 103 "
        "(General Chemistry II). Each course is organized into chapters and contains video lessons, "
        "downloadable notes, practice exercises, and quizzes. To enter a course, simply click on "
        "its card from the Courses page."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – OPENING A CHAPTER
# ══════════════════════════════════════════════════════════════════════════════
def slide07():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "CHAPTERS",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Opening a Chapter",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Flowchart
    cx = W / 2
    nodes = [
        ("🌐  Open chembridge.msnour.com", DEEP_PURPLE),
        ("📚  Click 'Courses' in the menu", DEEP_PURPLE),
        ("🎯  Select CHEM 101 or CHEM 103", DEEP_PURPLE),
        ("📖  Browse the Chapter List", LIGHT_PURPLE),
        ("▶️   Click on a Chapter to Open it", SOFT_PINK),
        ("📂  Access all Chapter Materials", RGBColor(0x2E, 0x86, 0x48)),
    ]
    node_w, node_h = Inches(4.5), Inches(0.52)
    start_y = Inches(1.62)
    spacing = Inches(0.82)
    for i, (label, bg) in enumerate(nodes):
        ty = start_y + i * spacing
        nx = cx - node_w / 2
        add_rect(slide, nx, ty, node_w, node_h,
                 fill_color=bg, line_color=WHITE, line_width=Pt(0.5))
        add_text_box(slide, label,
                     nx + Inches(0.15), ty + Inches(0.1),
                     node_w - Inches(0.3), node_h - Inches(0.12),
                     font_size=Pt(11), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            arrow_y = ty + node_h
            add_text_box(slide, "▼",
                         cx - Inches(0.18), arrow_y + Inches(0.05),
                         Inches(0.36), Inches(0.2),
                         font_size=Pt(10), font_color=SOFT_PINK, align=PP_ALIGN.CENTER)

    # Chapter list mockup on right
    list_l = Inches(8.3)
    add_rect(slide, list_l, Inches(1.65), Inches(4.7), Inches(5.25),
             fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1.2))
    add_text_box(slide, "CHEM 101 – Chapter List",
                 list_l + Inches(0.2), Inches(1.78), Inches(4.3), Inches(0.38),
                 font_size=Pt(11), font_color=SOFT_PINK, bold=True)
    add_rect(slide, list_l + Inches(0.2), Inches(2.14), Inches(4.3), Inches(0.03),
             fill_color=SOFT_PINK)
    chapters = [
        "Chapter 1: Atomic Structure & Periodicity",
        "Chapter 2: Chemical Bonding",
        "Chapter 3: Molecular Geometry (VSEPR)",
        "Chapter 4: Thermodynamics",
        "Chapter 5: Chemical Kinetics",
        "Chapter 6: Equilibrium",
        "Chapter 7: Acids & Bases",
        "Chapter 8: Electrochemistry",
    ]
    for i, ch in enumerate(chapters):
        ty = Inches(2.22) + i * Inches(0.54)
        bg = RGBColor(0x55, 0x30, 0x80) if i % 2 == 0 else RGBColor(0x3E, 0x1D, 0x68)
        add_rect(slide, list_l + Inches(0.2), ty, Inches(4.3), Inches(0.5), fill_color=bg)
        add_text_box(slide, "▶  " + ch,
                     list_l + Inches(0.35), ty + Inches(0.1),
                     Inches(4.0), Inches(0.36),
                     font_size=Pt(8.5), font_color=WHITE)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "To open a chapter, start by going to Courses in the top navigation, select your course "
        "(CHEM 101 or CHEM 103), and you'll see a list of all available chapters. "
        "Simply click on any chapter to open it. Inside, you'll find video lessons, "
        "notes, exercises, and a quiz for that chapter."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – WATCHING LESSONS
# ══════════════════════════════════════════════════════════════════════════════
def slide08():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "VIDEO LESSONS",
                 Inches(0.5), Inches(0.22), Inches(2.8), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Watching Lessons",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Video player mockup
    vl, vt = Inches(0.5), Inches(1.58)
    vw, vh = Inches(7.8), Inches(4.45)
    add_rect(slide, vl, vt, vw, vh,
             fill_color=RGBColor(0x0A, 0x0A, 0x1A), line_color=SOFT_PINK, line_width=Pt(1.5))

    # Video content area
    add_rect(slide, vl + Inches(0.15), vt + Inches(0.15),
             vw - Inches(0.3), vh - Inches(0.95),
             fill_color=RGBColor(0x1A, 0x0A, 0x30))
    # Play button
    play_cx = vl + vw / 2
    play_cy = vt + (vh - Inches(0.95)) / 2 + Inches(0.15)
    add_circle(slide, play_cx, play_cy, Inches(0.55), SOFT_PINK)
    add_text_box(slide, "▶",
                 play_cx - Inches(0.55), play_cy - Inches(0.55),
                 Inches(1.1), Inches(1.1),
                 font_size=Pt(22), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Screen label
    add_text_box(slide, "Chapter 2: Chemical Bonding – Lecture Video",
                 vl + Inches(0.3), vt + Inches(0.25), Inches(5.0), Inches(0.35),
                 font_size=Pt(9), font_color=SOFT_PINK)

    # Control bar
    ctrl_t = vt + vh - Inches(0.82)
    add_rect(slide, vl, ctrl_t, vw, Inches(0.82),
             fill_color=RGBColor(0x22, 0x11, 0x44))
    # Progress bar
    add_rect(slide, vl + Inches(0.2), ctrl_t + Inches(0.12),
             vw - Inches(0.4), Inches(0.08),
             fill_color=RGBColor(0x55, 0x55, 0x66))
    add_rect(slide, vl + Inches(0.2), ctrl_t + Inches(0.12),
             Inches(3.0), Inches(0.08), fill_color=SOFT_PINK)
    # Buttons
    for icon, bx in [("⏮", vl + Inches(0.3)), ("⏸", vl + Inches(0.85)),
                     ("⏭", vl + Inches(1.4))]:
        add_text_box(slide, icon, bx, ctrl_t + Inches(0.3),
                     Inches(0.42), Inches(0.35), font_size=Pt(14), font_color=WHITE)
    add_text_box(slide, "0:00 / 18:45",
                 vl + Inches(2.1), ctrl_t + Inches(0.32),
                 Inches(1.5), Inches(0.3),
                 font_size=Pt(9), font_color=SOFT_PINK)
    add_text_box(slide, "🔊  1x  ⛶  ⚙️",
                 vl + vw - Inches(2.0), ctrl_t + Inches(0.32),
                 Inches(1.8), Inches(0.3),
                 font_size=Pt(9), font_color=WHITE)

    # Right panel
    rp_l = Inches(8.65)
    add_text_box(slide, "Video Controls Guide",
                 rp_l, Inches(1.62), Inches(4.4), Inches(0.45),
                 font_size=Pt(13), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, rp_l, Inches(2.05), Inches(3.0), Inches(0.04), fill_color=SOFT_PINK)

    controls = [
        ("▶ / ⏸", "Play or pause the video"),
        ("⏮ / ⏭",  "Skip to previous/next lecture"),
        ("Progress bar", "Click anywhere to jump to a time"),
        ("1x speed", "Change playback speed (0.75x – 2x)"),
        ("⛶ fullscreen", "Watch in full-screen mode"),
        ("⚙️ settings", "Change quality or subtitles"),
    ]
    for i, (ctrl, desc) in enumerate(controls):
        ty = Inches(2.2) + i * Inches(0.65)
        add_rect(slide, rp_l, ty, Inches(4.4), Inches(0.58),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.6))
        add_text_box(slide, ctrl,
                     rp_l + Inches(0.1), ty + Inches(0.06),
                     Inches(1.5), Inches(0.26),
                     font_size=Pt(9), font_color=DEEP_PURPLE, bold=True)
        add_text_box(slide, desc,
                     rp_l + Inches(0.1), ty + Inches(0.3),
                     Inches(4.1), Inches(0.26),
                     font_size=Pt(8.5), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "Inside each chapter you'll find video lectures recorded by Teacher Nourhan. "
        "The video player has full controls: play/pause, skip forward or back, a progress bar "
        "to jump to any point, playback speed control (great for review at 1.5x), "
        "and a fullscreen mode. You can re-watch videos as many times as you need."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – STUDY MATERIALS
# ══════════════════════════════════════════════════════════════════════════════
def slide09():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)
    add_molecule_decoration(slide, Inches(0.6), Inches(6.8), 0.18, PALE_PINK)

    add_text_box(slide, "RESOURCES",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Accessing Study Materials",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    materials = [
        ("📝", "Lecture Notes",
         "Comprehensive PDF notes for each chapter,\naligned with the video content.",
         ["Chapter summaries", "Key definitions", "Important equations", "Diagrams & charts"]),
        ("🧮", "Formula Sheets",
         "Quick-reference sheets with all formulas\nyou need for exams.",
         ["Organized by topic", "With units & constants", "Exam-ready format", "Printable PDF"]),
        ("📖", "Chapter Summaries",
         "Condensed review sheets highlighting\nthe most critical concepts.",
         ["Bullet-point format", "Topic overview", "Concept maps", "Key takeaways"]),
        ("📥", "Downloadable Files",
         "All materials available for offline study —\ndownload and study anywhere.",
         ["PDF format", "High quality print", "Mobile-friendly", "Free to download"]),
    ]

    for i, (icon, title, desc, bullets) in enumerate(materials):
        col = i % 2
        row = i // 2
        lx = Inches(0.5) + col * Inches(6.35)
        ty = Inches(1.78) + row * Inches(2.58)
        w, h = Inches(6.1), Inches(2.42)
        add_rect(slide, lx, ty, w, h,
                 fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1.2))
        add_text_box(slide, icon,
                     lx + Inches(0.18), ty + Inches(0.2), Inches(0.7), Inches(0.7),
                     font_size=Pt(28), font_color=SOFT_PINK)
        add_text_box(slide, title,
                     lx + Inches(0.95), ty + Inches(0.18), w - Inches(1.1), Inches(0.42),
                     font_size=Pt(13), font_color=WHITE, bold=True)
        add_text_box(slide, desc,
                     lx + Inches(0.95), ty + Inches(0.6), w - Inches(1.1), Inches(0.55),
                     font_size=Pt(9), font_color=SOFT_PINK)
        for j, b in enumerate(bullets):
            bx = lx + Inches(0.18) + (j % 2) * Inches(2.95)
            by = ty + Inches(1.28) + (j // 2) * Inches(0.45)
            add_text_box(slide, "✓  " + b, bx, by, Inches(2.8), Inches(0.38),
                         font_size=Pt(8.5), font_color=WHITE)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "ChemBridge provides four types of study materials for each chapter: "
        "detailed lecture notes, formula reference sheets, chapter summaries, and downloadable PDFs. "
        "Everything is available inside the chapter page. Click the download icon next to any "
        "material to save it to your device for offline study."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – PRACTICE QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════
def slide10():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "PRACTICE",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Solving Practice Questions",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Left – practice question mockup
    ql, qt = Inches(0.5), Inches(1.58)
    qw, qh = Inches(6.5), Inches(5.38)
    add_rect(slide, ql, qt, qw, qh,
             fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(1.5))
    add_text_box(slide, "📋  Practice Set – Chapter 2: Chemical Bonding",
                 ql + Inches(0.25), qt + Inches(0.2), qw - Inches(0.5), Inches(0.38),
                 font_size=Pt(10), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, ql + Inches(0.25), qt + Inches(0.58), qw - Inches(0.5), Inches(0.03),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Question 3 of 12",
                 ql + Inches(0.25), qt + Inches(0.68), Inches(2.5), Inches(0.3),
                 font_size=Pt(9), font_color=SOFT_PINK)
    # Progress bar
    add_rect(slide, ql + Inches(0.25), qt + Inches(0.98), qw - Inches(0.5), Inches(0.1),
             fill_color=RGBColor(0xCC, 0xCC, 0xDD))
    add_rect(slide, ql + Inches(0.25), qt + Inches(0.98), (qw - Inches(0.5)) * 3 / 12, Inches(0.1),
             fill_color=SOFT_PINK)

    add_text_box(slide,
                 "Which of the following molecules has a trigonal planar geometry?",
                 ql + Inches(0.25), qt + Inches(1.2), qw - Inches(0.5), Inches(0.55),
                 font_size=Pt(11), font_color=DEEP_PURPLE, bold=True)

    options = ["A)  NH₃  (Ammonia)", "B)  BF₃  (Boron Trifluoride)",
               "C)  H₂O  (Water)", "D)  CH₄  (Methane)"]
    correct = 1
    for i, opt in enumerate(options):
        oy = qt + Inches(1.88) + i * Inches(0.72)
        bg = RGBColor(0x2E, 0x86, 0x48) if i == correct else DEEP_PURPLE
        border = WHITE if i == correct else SOFT_PINK
        add_rect(slide, ql + Inches(0.25), oy, qw - Inches(0.5), Inches(0.58),
                 fill_color=bg, line_color=border, line_width=Pt(0.8))
        add_text_box(slide, opt,
                     ql + Inches(0.45), oy + Inches(0.12), qw - Inches(0.9), Inches(0.38),
                     font_size=Pt(10), font_color=WHITE)
        if i == correct:
            add_text_box(slide, "✓ Correct!",
                         ql + qw - Inches(1.4), oy + Inches(0.12), Inches(1.2), Inches(0.38),
                         font_size=Pt(9), font_color=WHITE, bold=True)

    # Explanation box
    exp_y = qt + Inches(4.78)
    add_rect(slide, ql + Inches(0.25), exp_y, qw - Inches(0.5), Inches(0.45),
             fill_color=DEEP_PURPLE)
    add_text_box(slide, "💡  BF₃ has 3 bonding pairs and no lone pairs → trigonal planar (120°)",
                 ql + Inches(0.35), exp_y + Inches(0.08), qw - Inches(0.7), Inches(0.35),
                 font_size=Pt(9), font_color=SOFT_PINK)

    # Right panel
    rp = Inches(7.4)
    add_text_box(slide, "Practice Features",
                 rp, Inches(1.62), Inches(5.5), Inches(0.45),
                 font_size=Pt(14), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, rp, Inches(2.05), Inches(3.2), Inches(0.04), fill_color=SOFT_PINK)

    feats = [
        ("📦", "Question Banks",     "Multiple questions per chapter topic"),
        ("✏️", "Multiple Choice",    "Select the best answer from 4 options"),
        ("💡", "Instant Feedback",   "See correct answer & explanation right away"),
        ("🔁", "Retry Questions",    "Redo missed questions until mastered"),
        ("📊", "Score Tracking",     "Each attempt is saved to your progress"),
        ("🔀", "Randomized Order",   "Questions shuffled to keep you sharp"),
    ]
    for i, (icon, title, desc) in enumerate(feats):
        ty = Inches(2.22) + i * Inches(0.78)
        add_rect(slide, rp, ty, Inches(5.55), Inches(0.68),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.6))
        add_text_box(slide, icon, rp + Inches(0.1), ty + Inches(0.1),
                     Inches(0.45), Inches(0.45), font_size=Pt(16))
        add_text_box(slide, title, rp + Inches(0.62), ty + Inches(0.08),
                     Inches(2.0), Inches(0.28), font_size=Pt(10), font_color=DEEP_PURPLE, bold=True)
        add_text_box(slide, desc, rp + Inches(0.62), ty + Inches(0.36),
                     Inches(4.7), Inches(0.28), font_size=Pt(8.5), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "After watching a video and reading the notes, it's time to practice. "
        "Each chapter has a practice question set with multiple-choice questions. "
        "When you select an answer, you immediately see whether it's correct and "
        "why — the explanation helps you learn from mistakes. You can retry any "
        "question as many times as needed."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – QUIZZES
# ══════════════════════════════════════════════════════════════════════════════
def slide11():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "ASSESSMENT",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Taking Chapter Quizzes",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Quiz workflow – horizontal
    steps = [
        ("1", "Open\nChapter", "Navigate to\nthe chapter"),
        ("2", "Click\n'Take Quiz'", "Find the Quiz\nbutton at bottom"),
        ("3", "Read\nInstructions", "Note time limit\n& question count"),
        ("4", "Answer\nAll Questions", "Select best\nanswer for each"),
        ("5", "Submit\nQuiz", "Click Submit\nwhen done"),
        ("6", "View\nResults", "See score &\ncorrect answers"),
    ]
    step_w = Inches(1.88)
    step_h = Inches(1.75)
    total_w = len(steps) * step_w + (len(steps) - 1) * Inches(0.22)
    sx = (W - total_w) / 2
    sy = Inches(1.65)
    for i, (num, title, body) in enumerate(steps):
        lx = sx + i * (step_w + Inches(0.22))
        add_rect(slide, lx, sy, step_w, step_h,
                 fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1.2))
        add_circle(slide, lx + step_w / 2, sy + Inches(0.35), Inches(0.28), SOFT_PINK)
        add_text_box(slide, num,
                     lx + step_w / 2 - Inches(0.28), sy + Inches(0.1),
                     Inches(0.56), Inches(0.5),
                     font_size=Pt(13), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, title,
                     lx + Inches(0.1), sy + Inches(0.72), step_w - Inches(0.2), Inches(0.55),
                     font_size=Pt(10), font_color=SOFT_PINK, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     lx + Inches(0.1), sy + Inches(1.27), step_w - Inches(0.2), Inches(0.45),
                     font_size=Pt(8.5), font_color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = lx + step_w + Inches(0.04)
            add_text_box(slide, "→",
                         ax, sy + step_h / 2 - Inches(0.15),
                         Inches(0.18), Inches(0.3),
                         font_size=Pt(14), font_color=SOFT_PINK, align=PP_ALIGN.CENTER)

    # Results card mockup
    rl, rt = Inches(0.5), Inches(3.65)
    rw, rh = Inches(5.8), Inches(3.25)
    add_rect(slide, rl, rt, rw, rh,
             fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1.5))
    add_text_box(slide, "📊  Quiz Results – Chapter 2: Chemical Bonding",
                 rl + Inches(0.25), rt + Inches(0.2), rw - Inches(0.5), Inches(0.38),
                 font_size=Pt(11), font_color=SOFT_PINK, bold=True)
    add_rect(slide, rl + Inches(0.25), rt + Inches(0.58), rw - Inches(0.5), Inches(0.03),
             fill_color=SOFT_PINK)
    # Score donut placeholder
    add_circle(slide, rl + Inches(1.3), rt + Inches(1.88), Inches(0.88),
               RGBColor(0x55, 0x30, 0x80))
    add_circle(slide, rl + Inches(1.3), rt + Inches(1.88), Inches(0.62), DEEP_PURPLE)
    add_text_box(slide, "85%",
                 rl + Inches(0.75), rt + Inches(1.62), Inches(1.1), Inches(0.52),
                 font_size=Pt(16), font_color=SOFT_PINK, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Score",
                 rl + Inches(0.75), rt + Inches(2.12), Inches(1.1), Inches(0.28),
                 font_size=Pt(8), font_color=WHITE, align=PP_ALIGN.CENTER)
    # Stats
    for label, val, cx2 in [("Correct", "17/20", rl + Inches(2.8)),
                             ("Time", "12 min", rl + Inches(4.0))]:
        add_text_box(slide, val,
                     cx2 - Inches(0.6), rt + Inches(1.5), Inches(1.2), Inches(0.55),
                     font_size=Pt(18), font_color=SOFT_PINK, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     cx2 - Inches(0.6), rt + Inches(2.02), Inches(1.2), Inches(0.3),
                     font_size=Pt(9), font_color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, rl + Inches(0.25), rt + Inches(2.6), rw - Inches(0.5), Inches(0.35),
             fill_color=SOFT_PINK)
    add_text_box(slide, "🎉  Great job! Review the 3 missed questions below.",
                 rl + Inches(0.35), rt + Inches(2.66), rw - Inches(0.7), Inches(0.28),
                 font_size=Pt(9), font_color=WHITE, bold=True)

    # Right: quiz tips
    tips_l = Inches(6.75)
    add_text_box(slide, "Quiz Tips for Success",
                 tips_l, Inches(3.68), Inches(6.2), Inches(0.42),
                 font_size=Pt(13), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, tips_l, Inches(4.08), Inches(3.5), Inches(0.04), fill_color=SOFT_PINK)

    qtips = [
        "📝  Read the question carefully before answering",
        "⏱️  Keep an eye on the timer — don't spend too long on one question",
        "🔄  You can review questions before submitting",
        "💡  After submission, study the explanations for wrong answers",
        "📈  Aim for 80%+ before moving to the next chapter",
    ]
    for i, tip in enumerate(qtips):
        ty = Inches(4.22) + i * Inches(0.55)
        add_rect(slide, tips_l, ty, Inches(6.2), Inches(0.48),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.6))
        add_text_box(slide, tip, tips_l + Inches(0.1), ty + Inches(0.1),
                     Inches(5.9), Inches(0.36), font_size=Pt(9.5), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "Quizzes are the best way to test your understanding before an exam. "
        "Open the chapter, scroll to the bottom, and click 'Take Quiz'. "
        "Answer all questions, then submit to see your score. "
        "The results page shows your percentage, how many you got right, "
        "and explanations for any missed questions. Aim for 80% or higher before moving on."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – TRACKING PROGRESS
# ══════════════════════════════════════════════════════════════════════════════
def slide12():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)

    add_text_box(slide, "ANALYTICS",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Tracking Your Progress",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Dashboard card top row
    stats = [
        ("8 / 12", "Lessons\nCompleted", "📹"),
        ("75%",    "Average\nQuiz Score", "🏆"),
        ("3",      "Chapters\nMastered",  "⭐"),
        ("4 hrs",  "Study\nTime",         "⏱️"),
    ]
    cw, ch = Inches(2.95), Inches(1.35)
    for i, (val, label, icon) in enumerate(stats):
        lx = Inches(0.5) + i * (cw + Inches(0.12))
        ty = Inches(1.62)
        add_rect(slide, lx, ty, cw, ch,
                 fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(1))
        add_text_box(slide, icon, lx + Inches(0.18), ty + Inches(0.18),
                     Inches(0.55), Inches(0.55), font_size=Pt(22))
        add_text_box(slide, val,
                     lx + Inches(0.78), ty + Inches(0.14), cw - Inches(0.88), Inches(0.55),
                     font_size=Pt(22), font_color=SOFT_PINK, bold=True)
        add_text_box(slide, label,
                     lx + Inches(0.78), ty + Inches(0.68), cw - Inches(0.88), Inches(0.55),
                     font_size=Pt(9), font_color=WHITE)

    # Chapter progress table
    table_l, table_t = Inches(0.5), Inches(3.15)
    table_w = Inches(7.3)
    add_rect(slide, table_l, table_t, table_w, Inches(0.42),
             fill_color=DEEP_PURPLE)
    for col, (hdr, cw2) in enumerate([("Chapter", Inches(2.8)), ("Status", Inches(1.4)),
                                      ("Score", Inches(1.2)), ("Completed", Inches(1.9))]):
        lx = table_l + sum(Inches(2.8) if x == 0 else Inches(1.4) if x == 1 else Inches(1.2) if x == 2 else Inches(1.9)
                           for x in range(col))
        lx = table_l + [0, Inches(2.8), Inches(4.2), Inches(5.4)][col]
        add_text_box(slide, hdr, lx + Inches(0.1), table_t + Inches(0.1),
                     cw2, Inches(0.28),
                     font_size=Pt(9), font_color=SOFT_PINK, bold=True)

    rows = [
        ("Chapter 1: Atomic Structure",   "✅ Complete", "90%", "May 12"),
        ("Chapter 2: Chemical Bonding",   "✅ Complete", "85%", "May 18"),
        ("Chapter 3: Molecular Geometry", "✅ Complete", "78%", "May 25"),
        ("Chapter 4: Thermodynamics",     "▶ In Progress","—",  "—"),
        ("Chapter 5: Kinetics",           "🔒 Locked",   "—",  "—"),
        ("Chapter 6: Equilibrium",        "🔒 Locked",   "—",  "—"),
    ]
    for i, (ch_name, status, score, date) in enumerate(rows):
        ry = table_t + Inches(0.42) + i * Inches(0.48)
        bg = PALE_PINK if i % 2 == 0 else WHITE
        add_rect(slide, table_l, ry, table_w, Inches(0.46), fill_color=bg)
        for ci, (txt, cx_off) in enumerate(
                [(ch_name, 0), (status, Inches(2.8)),
                 (score, Inches(4.2)), (date, Inches(5.4))]):
            fc = RGBColor(0x2E, 0x86, 0x48) if "✅" in txt else (
                 SOFT_PINK if "▶" in txt else DARK_GRAY)
            add_text_box(slide, txt,
                         table_l + cx_off + Inches(0.1), ry + Inches(0.1),
                         Inches(2.6), Inches(0.32),
                         font_size=Pt(8.5), font_color=fc)

    # Right: performance tips
    rp = Inches(8.2)
    add_text_box(slide, "Progress Indicators",
                 rp, Inches(3.18), Inches(4.8), Inches(0.42),
                 font_size=Pt(13), font_color=DEEP_PURPLE, bold=True)
    add_rect(slide, rp, Inches(3.58), Inches(3.2), Inches(0.04), fill_color=SOFT_PINK)

    indicators = [
        ("✅", "Complete",   "All lessons & quiz done (>80%)"),
        ("▶",  "In Progress","Currently working through chapter"),
        ("🔒", "Locked",     "Complete previous chapter first"),
        ("⭐", "Mastered",   "Quiz score ≥ 90% — excellent!"),
    ]
    for i, (icon, label, desc) in enumerate(indicators):
        ty = Inches(3.72) + i * Inches(0.78)
        add_rect(slide, rp, ty, Inches(4.8), Inches(0.68),
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(0.6))
        add_text_box(slide, icon, rp + Inches(0.1), ty + Inches(0.1),
                     Inches(0.38), Inches(0.45), font_size=Pt(16), align=PP_ALIGN.CENTER)
        add_text_box(slide, label, rp + Inches(0.6), ty + Inches(0.08),
                     Inches(4.0), Inches(0.28), font_size=Pt(10), font_color=DEEP_PURPLE, bold=True)
        add_text_box(slide, desc, rp + Inches(0.6), ty + Inches(0.35),
                     Inches(4.0), Inches(0.28), font_size=Pt(8.5), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "The Progress page gives you a dashboard view of your learning journey. "
        "You can see how many lessons you've completed, your average quiz score, "
        "how many chapters you've mastered, and your total study time. "
        "The chapter table shows the status of each chapter — complete, in progress, or locked. "
        "Focus on completing each chapter before moving to the next."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – RECOMMENDED STUDY WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
def slide13():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_molecule_decoration(slide, Inches(0.7), Inches(6.5), 0.22, PALE_PINK)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)
    add_atom_ring(slide, W - Inches(1.0), Inches(6.5), Inches(0.45), PALE_PINK)

    add_text_box(slide, "STUDY GUIDE",
                 Inches(0.5), Inches(0.22), Inches(2.5), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Recommended Study Workflow",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Central spine
    spine_x = W / 2 - Inches(0.04)
    spine_top = Inches(1.65)
    spine_bot = Inches(6.95)
    add_rect(slide, spine_x, spine_top, Inches(0.08),
             spine_bot - spine_top, fill_color=SOFT_PINK)

    steps = [
        ("📹", "Watch the Lecture",
         "Open the chapter and watch the video lesson. Take notes in your notebook.",
         True),
        ("📝", "Read the Notes",
         "Download and read the lecture notes PDF. Highlight key formulas and definitions.",
         False),
        ("✏️", "Solve Practice Questions",
         "Work through the practice set. Don't skip — even difficult questions!",
         True),
        ("🧪", "Take the Chapter Quiz",
         "Test yourself with the timed quiz. Aim for 80% or higher.",
         False),
        ("🔍", "Review Mistakes",
         "Go back to any questions you got wrong. Re-read the explanation.",
         True),
        ("⭐", "Master the Chapter",
         "Once you score 80%+ on the quiz, you're ready for the next chapter!",
         False),
    ]
    node_w = Inches(5.4)
    node_h = Inches(0.72)
    step_gap = Inches(0.86)
    for i, (icon, title, body, left_side) in enumerate(steps):
        ty = Inches(1.65) + i * step_gap
        if left_side:
            lx = spine_x - node_w - Inches(0.35)
        else:
            lx = spine_x + Inches(0.35)
        bg = DEEP_PURPLE if i % 3 != 2 else LIGHT_PURPLE
        add_rect(slide, lx, ty, node_w, node_h,
                 fill_color=bg, line_color=SOFT_PINK, line_width=Pt(1))
        # icon circle
        ic_cx = lx + (Inches(0.38) if left_side else node_w - Inches(0.38))
        add_circle(slide, ic_cx, ty + node_h / 2, Inches(0.28), SOFT_PINK)
        add_text_box(slide, icon,
                     ic_cx - Inches(0.28), ty + node_h / 2 - Inches(0.28),
                     Inches(0.56), Inches(0.56),
                     font_size=Pt(14), align=PP_ALIGN.CENTER)
        tx = lx + Inches(0.7) if left_side else lx + Inches(0.1)
        add_text_box(slide, title,
                     tx, ty + Inches(0.06), node_w - Inches(0.8), Inches(0.28),
                     font_size=Pt(10), font_color=SOFT_PINK, bold=True)
        add_text_box(slide, body,
                     tx, ty + Inches(0.35), node_w - Inches(0.8), Inches(0.34),
                     font_size=Pt(8.5), font_color=WHITE)
        # dot on spine
        add_circle(slide, spine_x + Inches(0.04), ty + node_h / 2, Inches(0.1), SOFT_PINK)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "Here's the recommended way to study each chapter on ChemBridge. "
        "Start by watching the video lecture, then read the notes PDF. "
        "Next, work through the practice questions and take the chapter quiz. "
        "If you make mistakes, review them using the explanations. "
        "Once you score 80% or higher on the quiz, you've mastered the chapter "
        "and can confidently move on to the next one."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CONTACT AND SUPPORT
# ══════════════════════════════════════════════════════════════════════════════
def slide14():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, WHITE)
    gradient_header(slide)
    add_atom_ring(slide, W - Inches(0.9), Inches(0.7), Inches(0.55), SOFT_PINK)
    add_molecule_decoration(slide, Inches(12.5), Inches(6.2), 0.22, PALE_PINK)

    add_text_box(slide, "SUPPORT",
                 Inches(0.5), Inches(0.22), Inches(2.0), Inches(0.5),
                 font_size=Pt(9), font_color=SOFT_PINK, bold=True)
    add_text_box(slide, "Contact & Support",
                 Inches(0.5), Inches(0.52), Inches(10.0), Inches(0.7),
                 font_size=Pt(28), font_color=WHITE, bold=True)

    # Teacher profile card (left)
    pc_l, pc_t = Inches(0.5), Inches(1.62)
    pc_w, pc_h = Inches(4.8), Inches(5.2)
    add_rect(slide, pc_l, pc_t, pc_w, pc_h,
             fill_color=DEEP_PURPLE, line_color=SOFT_PINK, line_width=Pt(2))
    # Avatar
    av_cx = pc_l + pc_w / 2
    av_cy = pc_t + Inches(1.0)
    add_circle(slide, av_cx, av_cy, Inches(0.75), SOFT_PINK)
    add_circle(slide, av_cx, av_cy, Inches(0.65), LIGHT_PURPLE)
    add_text_box(slide, "NZ",
                 av_cx - Inches(0.65), av_cy - Inches(0.65),
                 Inches(1.3), Inches(1.3),
                 font_size=Pt(22), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Teacher Nourhan Zaher",
                 pc_l + Inches(0.2), pc_t + Inches(2.0), pc_w - Inches(0.4), Inches(0.52),
                 font_size=Pt(14), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Chemistry Instructor\nQatar University – CHEM 101 & CHEM 103",
                 pc_l + Inches(0.2), pc_t + Inches(2.52), pc_w - Inches(0.4), Inches(0.7),
                 font_size=Pt(10), font_color=SOFT_PINK, align=PP_ALIGN.CENTER)
    add_rect(slide, pc_l + Inches(0.5), pc_t + Inches(3.32), pc_w - Inches(1.0), Inches(0.04),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Available for:\nQuestions · Support · Guidance",
                 pc_l + Inches(0.2), pc_t + Inches(3.48), pc_w - Inches(0.4), Inches(0.58),
                 font_size=Pt(10), font_color=WHITE, align=PP_ALIGN.CENTER)
    # Hours badge
    add_rect(slide, pc_l + Inches(0.5), pc_t + Inches(4.22), pc_w - Inches(1.0), Inches(0.42),
             fill_color=SOFT_PINK)
    add_text_box(slide, "Responds within 24 hours",
                 pc_l + Inches(0.5), pc_t + Inches(4.27), pc_w - Inches(1.0), Inches(0.34),
                 font_size=Pt(9.5), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Contact methods
    contacts = [
        ("📱", "WhatsApp",
         "33625942",
         "Send a message anytime for quick help with questions or content access."),
        ("✉️", "Email",
         "noorzaher164@gmail.com",
         "For detailed questions, material requests, or account issues."),
        ("🌐", "Website",
         "chembridge.msnour.com",
         "Your main learning hub — always available 24/7."),
    ]
    for i, (icon, label, contact, desc) in enumerate(contacts):
        cl = Inches(6.0)
        ct = Inches(1.68) + i * Inches(1.75)
        cw2, ch2 = Inches(7.0), Inches(1.58)
        add_rect(slide, cl, ct, cw2, ch2,
                 fill_color=PALE_PINK, line_color=SOFT_PINK, line_width=Pt(1.2))
        # icon circle
        add_circle(slide, cl + Inches(0.62), ct + ch2 / 2, Inches(0.42), DEEP_PURPLE)
        add_text_box(slide, icon,
                     cl + Inches(0.2), ct + ch2 / 2 - Inches(0.42),
                     Inches(0.84), Inches(0.84),
                     font_size=Pt(20), align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     cl + Inches(1.3), ct + Inches(0.22), Inches(2.0), Inches(0.35),
                     font_size=Pt(10), font_color=SOFT_PINK, bold=True)
        add_text_box(slide, contact,
                     cl + Inches(1.3), ct + Inches(0.55), Inches(5.4), Inches(0.42),
                     font_size=Pt(15), font_color=DEEP_PURPLE, bold=True)
        add_text_box(slide, desc,
                     cl + Inches(1.3), ct + Inches(1.0), Inches(5.5), Inches(0.48),
                     font_size=Pt(9), font_color=DARK_GRAY)

    slide_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "If you ever need help, Teacher Nourhan is available through three channels. "
        "WhatsApp at 33625942 is the fastest way to get a quick answer. "
        "For longer questions, email noorzaher164@gmail.com. "
        "And of course, the platform itself at chembridge.msnour.com is available 24/7. "
        "Don't hesitate to ask — that's what ChemBridge is here for!"
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
def slide15():
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, DEEP_PURPLE)

    # Background chemistry decorations
    for cx, cy, r in [
        (Inches(1.0), Inches(1.5), Inches(0.8)),
        (Inches(12.0), Inches(1.8), Inches(0.6)),
        (Inches(0.8), Inches(6.5), Inches(0.5)),
        (Inches(12.5), Inches(6.0), Inches(0.7)),
        (Inches(6.5), Inches(0.6), Inches(0.4)),
    ]:
        add_atom_ring(slide, cx, cy, r, SOFT_PINK)

    add_molecule_decoration(slide, Inches(2.0), Inches(2.8), 0.28, SOFT_PINK)
    add_molecule_decoration(slide, Inches(11.0), Inches(5.5), 0.28, SOFT_PINK)
    add_molecule_decoration(slide, Inches(0.7), Inches(4.5), 0.2, LIGHT_PURPLE)
    add_molecule_decoration(slide, Inches(12.6), Inches(3.2), 0.2, LIGHT_PURPLE)

    # Soft radial glow (layered circles)
    for r_in, alpha in [(Inches(3.5), RGBColor(0x55, 0x30, 0x80)),
                        (Inches(2.8), RGBColor(0x4A, 0x25, 0x72))]:
        add_circle(slide, W / 2, H / 2, r_in, alpha)

    # Large logo circle
    add_circle(slide, W / 2, H / 2 - Inches(0.8), Inches(1.1), SOFT_PINK)
    add_circle(slide, W / 2, H / 2 - Inches(0.8), Inches(0.95), WHITE)
    add_text_box(slide, "CB",
                 W / 2 - Inches(0.95), H / 2 - Inches(0.8) - Inches(0.95),
                 Inches(1.9), Inches(1.9),
                 font_size=Pt(32), font_color=DEEP_PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, "Start Your Chemistry Journey Today",
                 Inches(1.5), Inches(3.55), W - Inches(3.0), Inches(0.8),
                 font_size=Pt(30), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, Inches(3.5), Inches(4.42), Inches(6.3), Inches(0.06),
             fill_color=SOFT_PINK)

    # Subtitle
    add_text_box(slide, "ChemBridge – Bridging the Gap Between You and Chemistry",
                 Inches(1.5), Inches(4.58), W - Inches(3.0), Inches(0.52),
                 font_size=Pt(14), font_color=SOFT_PINK, italic=True, align=PP_ALIGN.CENTER)

    # CTA button
    add_rect(slide, W / 2 - Inches(2.2), Inches(5.22), Inches(4.4), Inches(0.58),
             fill_color=SOFT_PINK)
    add_text_box(slide, "🌐  Visit chembridge.msnour.com  →",
                 W / 2 - Inches(2.2), Inches(5.28), Inches(4.4), Inches(0.48),
                 font_size=Pt(12), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Teacher credit
    add_text_box(slide, "by Teacher Nourhan Zaher",
                 Inches(2.0), Inches(6.0), W - Inches(4.0), Inches(0.35),
                 font_size=Pt(10), font_color=ACCENT_GOLD, italic=True, align=PP_ALIGN.CENTER)

    # Bottom bar
    add_rect(slide, 0, H - Inches(0.5), W, Inches(0.5), fill_color=SOFT_PINK)
    add_text_box(slide, "ChemBridge  |  CHEM 101 & CHEM 103  |  Qatar University  "
                 "|  chembridge.msnour.com  |  33625942",
                 Inches(0.3), H - Inches(0.44), W - Inches(0.6), Inches(0.38),
                 font_size=Pt(8.5), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "Thank you for your attention! ChemBridge is your complete chemistry learning companion "
        "for CHEM 101 and CHEM 103 at Qatar University. Visit chembridge.msnour.com to start "
        "your journey today. If you have any questions, contact Teacher Nourhan on WhatsApp at "
        "33625942 or email noorzaher164@gmail.com. You've got this — chemistry made easy!"
    )


# ══════════════════════════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ══════════════════════════════════════════════════════════════════════════════
print("Building slides...")
slide01()
print("  Slide 1/15 done")
slide02()
print("  Slide 2/15 done")
slide03()
print("  Slide 3/15 done")
slide04()
print("  Slide 4/15 done")
slide05()
print("  Slide 5/15 done")
slide06()
print("  Slide 6/15 done")
slide07()
print("  Slide 7/15 done")
slide08()
print("  Slide 8/15 done")
slide09()
print("  Slide 9/15 done")
slide10()
print("  Slide 10/15 done")
slide11()
print("  Slide 11/15 done")
slide12()
print("  Slide 12/15 done")
slide13()
print("  Slide 13/15 done")
slide14()
print("  Slide 14/15 done")
slide15()
print("  Slide 15/15 done")

out_path = "/home/user/Pearl-Vita-C-/ChemBridge_Student_Guide.pptx"
prs.save(out_path)
print(f"\n✅  Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
